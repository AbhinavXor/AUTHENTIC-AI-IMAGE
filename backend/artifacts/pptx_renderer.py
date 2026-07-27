from __future__ import annotations

from pathlib import Path
from tempfile import TemporaryDirectory
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from artifacts.charting import render_chart_image
from artifacts.models import (
    ArtifactDocument,
    ArtifactSection,
    BulletListBlock,
    ChartBlock,
    CodeBlock,
    ParagraphBlock,
    TableBlock,
)


SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

NAVY = RGBColor(23, 32, 51)
SLATE = RGBColor(85, 96, 114)
BODY = RGBColor(36, 43, 56)
MUTED = RGBColor(105, 115, 134)
LIGHT = RGBColor(233, 237, 244)
BORDER = RGBColor(205, 211, 221)
WHITE = RGBColor(255, 255, 255)
CODE_BG = RGBColor(245, 247, 250)


def _set_slide_background(
    slide,
    color: RGBColor = WHITE,
) -> None:
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_top_rule(
    slide,
) -> None:
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        0,
        SLIDE_WIDTH,
        Inches(0.08),
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = NAVY
    rule.line.fill.background()


def _add_footer(
    slide,
    *,
    slide_number: int,
    document_title: str,
) -> None:
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.55),
        Inches(7.08),
        Inches(12.23),
        Inches(0.015),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BORDER
    line.line.fill.background()

    left = slide.shapes.add_textbox(
        Inches(0.58),
        Inches(7.12),
        Inches(8.5),
        Inches(0.22),
    )
    paragraph = left.text_frame.paragraphs[0]
    paragraph.text = document_title[:80]
    paragraph.font.name = "Aptos"
    paragraph.font.size = Pt(8)
    paragraph.font.color.rgb = MUTED

    right = slide.shapes.add_textbox(
        Inches(11.85),
        Inches(7.12),
        Inches(0.9),
        Inches(0.22),
    )
    paragraph = right.text_frame.paragraphs[0]
    paragraph.text = str(slide_number)
    paragraph.alignment = PP_ALIGN.RIGHT
    paragraph.font.name = "Aptos"
    paragraph.font.size = Pt(8)
    paragraph.font.color.rgb = MUTED


def _add_slide_title(
    slide,
    title: str,
    *,
    subtitle: str | None = None,
) -> None:
    box = slide.shapes.add_textbox(
        Inches(0.65),
        Inches(0.42),
        Inches(12.0),
        Inches(0.7),
    )
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0

    paragraph = frame.paragraphs[0]
    paragraph.text = title
    paragraph.font.name = "Aptos Display"
    paragraph.font.size = Pt(25)
    paragraph.font.bold = True
    paragraph.font.color.rgb = NAVY

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.67),
            Inches(1.12),
            Inches(11.8),
            Inches(0.38),
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.clear()
        subtitle_frame.margin_left = 0
        subtitle_frame.margin_right = 0
        subtitle_frame.margin_top = 0
        subtitle_frame.margin_bottom = 0

        paragraph = subtitle_frame.paragraphs[0]
        paragraph.text = subtitle
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(11)
        paragraph.font.color.rgb = SLATE


def _add_title_slide(
    presentation: Presentation,
    artifact: ArtifactDocument,
) -> None:
    slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )
    _set_slide_background(slide, NAVY)

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.72),
        Inches(1.1),
        Inches(0.12),
        Inches(4.8),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(
        170,
        179,
        194,
    )
    accent.line.fill.background()

    title = slide.shapes.add_textbox(
        Inches(1.2),
        Inches(1.25),
        Inches(10.9),
        Inches(2.1),
    )
    frame = title.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    paragraph = frame.paragraphs[0]
    paragraph.text = artifact.title
    paragraph.font.name = "Aptos Display"
    paragraph.font.size = Pt(31)
    paragraph.font.bold = True
    paragraph.font.color.rgb = WHITE

    if artifact.subtitle:
        subtitle = slide.shapes.add_textbox(
            Inches(1.25),
            Inches(3.55),
            Inches(10.4),
            Inches(1.0),
        )
        frame = subtitle.text_frame
        frame.clear()
        frame.word_wrap = True

        paragraph = frame.paragraphs[0]
        paragraph.text = artifact.subtitle
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(17)
        paragraph.font.color.rgb = RGBColor(
            210,
            217,
            227,
        )

    author = slide.shapes.add_textbox(
        Inches(1.25),
        Inches(5.65),
        Inches(9.0),
        Inches(0.5),
    )
    paragraph = author.text_frame.paragraphs[0]
    paragraph.text = artifact.author or "Authentic AI"
    paragraph.font.name = "Aptos"
    paragraph.font.size = Pt(11)
    paragraph.font.color.rgb = RGBColor(
        190,
        199,
        213,
    )


def _add_agenda_slide(
    presentation: Presentation,
    artifact: ArtifactDocument,
    *,
    slide_number: int,
) -> None:
    slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )
    _set_slide_background(slide)
    _add_top_rule(slide)
    _add_slide_title(
        slide,
        "Contents",
        subtitle="Presentation structure",
    )

    box = slide.shapes.add_textbox(
        Inches(0.9),
        Inches(1.75),
        Inches(11.4),
        Inches(4.9),
    )
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True

    for index, section in enumerate(
        artifact.sections,
        start=1,
    ):
        paragraph = (
            frame.paragraphs[0]
            if index == 1
            else frame.add_paragraph()
        )
        paragraph.text = (
            f"{index:02d}  {section.title}"
        )
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = BODY
        paragraph.space_after = Pt(9)

    _add_footer(
        slide,
        slide_number=slide_number,
        document_title=artifact.title,
    )


def _add_section_divider(
    presentation: Presentation,
    artifact: ArtifactDocument,
    section: ArtifactSection,
    *,
    section_number: int,
    slide_number: int,
) -> None:
    slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )
    _set_slide_background(slide, LIGHT)

    number_box = slide.shapes.add_textbox(
        Inches(0.85),
        Inches(1.25),
        Inches(2.0),
        Inches(1.0),
    )
    paragraph = number_box.text_frame.paragraphs[0]
    paragraph.text = f"{section_number:02d}"
    paragraph.font.name = "Aptos Display"
    paragraph.font.size = Pt(46)
    paragraph.font.bold = True
    paragraph.font.color.rgb = SLATE

    title_box = slide.shapes.add_textbox(
        Inches(0.9),
        Inches(2.45),
        Inches(11.5),
        Inches(2.2),
    )
    frame = title_box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    paragraph = frame.paragraphs[0]
    paragraph.text = section.title
    paragraph.font.name = "Aptos Display"
    paragraph.font.size = Pt(30)
    paragraph.font.bold = True
    paragraph.font.color.rgb = NAVY

    _add_footer(
        slide,
        slide_number=slide_number,
        document_title=artifact.title,
    )


def _split_text(
    text: str,
    *,
    maximum_characters: int = 850,
) -> tuple[str, ...]:
    cleaned = " ".join(text.split())

    if len(cleaned) <= maximum_characters:
        return (cleaned,)

    chunks: list[str] = []
    remaining = cleaned

    while remaining:
        if len(remaining) <= maximum_characters:
            chunks.append(remaining)
            break

        split_at = remaining.rfind(
            ". ",
            0,
            maximum_characters,
        )

        if split_at < maximum_characters // 2:
            split_at = remaining.rfind(
                " ",
                0,
                maximum_characters,
            )

        if split_at <= 0:
            split_at = maximum_characters

        chunks.append(
            remaining[:split_at].strip()
        )
        remaining = remaining[split_at:].strip()

    return tuple(chunks)


def _add_paragraph_slide(
    presentation: Presentation,
    artifact: ArtifactDocument,
    section_title: str,
    text: str,
    *,
    continuation: int,
    slide_number: int,
) -> None:
    slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )
    _set_slide_background(slide)
    _add_top_rule(slide)

    title = section_title

    if continuation > 1:
        title = (
            f"{section_title} "
            f"({continuation})"
        )

    _add_slide_title(slide, title)

    body_box = slide.shapes.add_textbox(
        Inches(0.85),
        Inches(1.55),
        Inches(11.65),
        Inches(4.95),
    )
    frame = body_box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.12)
    frame.margin_right = Inches(0.12)
    frame.margin_top = Inches(0.08)
    frame.margin_bottom = Inches(0.08)
    frame.vertical_anchor = MSO_ANCHOR.TOP

    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.name = "Aptos"
    paragraph.font.size = Pt(18)
    paragraph.font.color.rgb = BODY
    paragraph.line_spacing = 1.15

    _add_footer(
        slide,
        slide_number=slide_number,
        document_title=artifact.title,
    )


def _chunk_items(
    items: Iterable[str],
    *,
    size: int = 6,
) -> tuple[tuple[str, ...], ...]:
    values = tuple(items)

    return tuple(
        values[index:index + size]
        for index in range(
            0,
            len(values),
            size,
        )
    )


def _add_list_slide(
    presentation: Presentation,
    artifact: ArtifactDocument,
    section_title: str,
    block: BulletListBlock,
    *,
    continuation: int,
    slide_number: int,
    items: tuple[str, ...],
) -> None:
    slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )
    _set_slide_background(slide)
    _add_top_rule(slide)

    title = section_title

    if continuation > 1:
        title = (
            f"{section_title} "
            f"({continuation})"
        )

    _add_slide_title(slide, title)

    body_box = slide.shapes.add_textbox(
        Inches(0.9),
        Inches(1.52),
        Inches(11.5),
        Inches(5.25),
    )
    frame = body_box.text_frame
    frame.clear()
    frame.word_wrap = True

    for index, item in enumerate(
        items,
        start=1,
    ):
        paragraph = (
            frame.paragraphs[0]
            if index == 1
            else frame.add_paragraph()
        )
        paragraph.text = (
            f"{index}. {item}"
            if block.ordered
            else item
        )
        paragraph.level = 0
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = BODY
        paragraph.space_after = Pt(10)
        paragraph.line_spacing = 1.08

        if not block.ordered:
            paragraph.text = f"•  {item}"

    _add_footer(
        slide,
        slide_number=slide_number,
        document_title=artifact.title,
    )


def _add_table_slide(
    presentation: Presentation,
    artifact: ArtifactDocument,
    section_title: str,
    block: TableBlock,
    *,
    slide_number: int,
) -> None:
    slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )
    _set_slide_background(slide)
    _add_top_rule(slide)
    _add_slide_title(
        slide,
        block.caption or section_title,
    )

    row_count = len(block.rows) + 1
    column_count = len(block.columns)

    table_shape = slide.shapes.add_table(
        row_count,
        column_count,
        Inches(0.65),
        Inches(1.55),
        Inches(12.05),
        Inches(4.95),
    )

    table = table_shape.table

    for column in table.columns:
        column.width = Inches(
            12.05 / column_count
        )

    for column_index, value in enumerate(
        block.columns
    ):
        cell = table.cell(
            0,
            column_index,
        )
        cell.text = value
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT

        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(11)
        paragraph.font.bold = True
        paragraph.font.color.rgb = NAVY

    for row_index, row_values in enumerate(
        block.rows,
        start=1,
    ):
        for column_index, value in enumerate(
            row_values
        ):
            cell = table.cell(
                row_index,
                column_index,
            )
            cell.text = value
            cell.vertical_anchor = (
                MSO_ANCHOR.MIDDLE
            )

            paragraph = (
                cell.text_frame.paragraphs[0]
            )
            paragraph.font.name = "Aptos"
            paragraph.font.size = Pt(10)
            paragraph.font.color.rgb = BODY

    _add_footer(
        slide,
        slide_number=slide_number,
        document_title=artifact.title,
    )


def _add_chart_slide(
    presentation: Presentation,
    artifact: ArtifactDocument,
    section_title: str,
    block: ChartBlock,
    *,
    chart_path: Path,
    slide_number: int,
) -> None:
    slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )
    _set_slide_background(slide)
    _add_top_rule(slide)
    _add_slide_title(
        slide,
        block.title or section_title,
        subtitle=block.caption,
    )

    slide.shapes.add_picture(
        str(chart_path),
        Inches(0.95),
        Inches(1.55),
        width=Inches(11.45),
        height=Inches(4.95),
    )

    _add_footer(
        slide,
        slide_number=slide_number,
        document_title=artifact.title,
    )


def _add_code_slide(
    presentation: Presentation,
    artifact: ArtifactDocument,
    section_title: str,
    block: CodeBlock,
    *,
    slide_number: int,
) -> None:
    slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )
    _set_slide_background(slide)
    _add_top_rule(slide)
    _add_slide_title(
        slide,
        section_title,
        subtitle=(
            block.language.upper()
            if block.language
            else "CODE"
        ),
    )

    background = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8),
        Inches(1.55),
        Inches(11.7),
        Inches(4.95),
    )
    background.fill.solid()
    background.fill.fore_color.rgb = CODE_BG
    background.line.color.rgb = BORDER

    code_box = slide.shapes.add_textbox(
        Inches(1.0),
        Inches(1.78),
        Inches(11.25),
        Inches(4.5),
    )
    frame = code_box.text_frame
    frame.clear()
    frame.word_wrap = True

    paragraph = frame.paragraphs[0]
    paragraph.text = block.code
    paragraph.font.name = "Courier New"
    paragraph.font.size = Pt(13)
    paragraph.font.color.rgb = BODY

    _add_footer(
        slide,
        slide_number=slide_number,
        document_title=artifact.title,
    )


def render_pptx(
    artifact: ArtifactDocument,
    output_path: Path,
) -> Path:
    output_path.parent.mkdir(
        parents=True,
        exist_ok=True,
    )

    presentation = Presentation()
    presentation.slide_width = SLIDE_WIDTH
    presentation.slide_height = SLIDE_HEIGHT

    presentation.core_properties.title = (
        artifact.title
    )
    presentation.core_properties.author = (
        artifact.author
        or "Authentic AI"
    )
    presentation.core_properties.subject = (
        artifact.subtitle
        or "Professional generated presentation"
    )
    presentation.core_properties.comments = (
        "Generated by Authentic AI Artifact Studio"
    )

    _add_title_slide(
        presentation,
        artifact,
    )

    slide_number = 2

    _add_agenda_slide(
        presentation,
        artifact,
        slide_number=slide_number,
    )
    slide_number += 1

    with TemporaryDirectory(
        prefix="authentic-pptx-charts-",
    ) as temporary_directory:
        chart_directory = Path(
            temporary_directory
        )
        chart_counter = 0

        for section_number, section in enumerate(
            artifact.sections,
            start=1,
        ):
            _add_section_divider(
                presentation,
                artifact,
                section,
                section_number=section_number,
                slide_number=slide_number,
            )
            slide_number += 1

            for block in section.blocks:
                if isinstance(
                    block,
                    ParagraphBlock,
                ):
                    for continuation, text in enumerate(
                        _split_text(block.text),
                        start=1,
                    ):
                        _add_paragraph_slide(
                            presentation,
                            artifact,
                            section.title,
                            text,
                            continuation=continuation,
                            slide_number=slide_number,
                        )
                        slide_number += 1

                elif isinstance(
                    block,
                    BulletListBlock,
                ):
                    for continuation, items in enumerate(
                        _chunk_items(
                            block.items
                        ),
                        start=1,
                    ):
                        _add_list_slide(
                            presentation,
                            artifact,
                            section.title,
                            block,
                            continuation=continuation,
                            slide_number=slide_number,
                            items=items,
                        )
                        slide_number += 1

                elif isinstance(
                    block,
                    TableBlock,
                ):
                    _add_table_slide(
                        presentation,
                        artifact,
                        section.title,
                        block,
                        slide_number=slide_number,
                    )
                    slide_number += 1

                elif isinstance(
                    block,
                    ChartBlock,
                ):
                    chart_counter += 1

                    chart_path = (
                        chart_directory
                        / f"chart-{chart_counter}.png"
                    )

                    render_chart_image(
                        block,
                        chart_path,
                    )

                    _add_chart_slide(
                        presentation,
                        artifact,
                        section.title,
                        block,
                        chart_path=chart_path,
                        slide_number=slide_number,
                    )
                    slide_number += 1

                elif isinstance(
                    block,
                    CodeBlock,
                ):
                    _add_code_slide(
                        presentation,
                        artifact,
                        section.title,
                        block,
                        slide_number=slide_number,
                    )
                    slide_number += 1

        presentation.save(output_path)

    if (
        not output_path.exists()
        or output_path.stat().st_size <= 0
    ):
        raise RuntimeError(
            "PPTX rendering failed."
        )

    return output_path
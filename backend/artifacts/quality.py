from __future__ import annotations

import html as html_lib
import io
import json
import re
from zipfile import BadZipFile, ZipFile
from dataclasses import dataclass, field, replace
from pathlib import Path
from typing import Any, Iterable

from artifacts.models import (
    ArtifactDocument,
    ArtifactSection,
    BulletListBlock,
    CalloutBlock,
    ChartBlock,
    ChartSeries,
    CodeBlock,
    DiagramBlock,
    EquationBlock,
    PageBreakBlock,
    ParagraphBlock,
    QuoteBlock,
    TableBlock,
)
from artifacts.visualization_blocks import count_authentic_chart_blocks
from core.artifact_settings import artifact_settings

_INTERNAL_CONTEXT = re.compile(
    r"<!--AUTHENTIC_[A-Z0-9_]+:[\s\S]*?-->"
)
_OUTER_FENCE = re.compile(
    r"^```(?:md|markdown)?\s*\n([\s\S]*?)\n```\s*$",
    re.IGNORECASE,
)
_HORIZONTAL_RULE = re.compile(
    r"^\s*(?:-{3,}|_{3,}|\*{3,})\s*$"
)
_MARKDOWN_TOKEN = re.compile(
    r"(?<!\\)(?:\*\*|__|~~|`)(.+?)(?<!\\)(?:\*\*|__|~~|`)"
)
_MARKDOWN_LINK = re.compile(r"\[([^\]]+)\]\(([^)]+)\)")
_INLINE_MATH = re.compile(r"(?<!\$)\$(?!\$)(.+?)(?<!\$)\$(?!\$)")
_RAW_MARKER = re.compile(
    r"(?:\*\*|~~|^\s*---\s*$)",
    re.MULTILINE,
)
_RAW_HTML = re.compile(
    r"<\/?(?:html|body|div|span|table|tr|td|script|style)\b",
    re.IGNORECASE,
)
_MARKDOWN_FENCE = re.compile(
    r"^```(?:md|markdown)[ \t]*\n([\s\S]*?)\n```[ \t]*(?=\n|\Z)",
    re.IGNORECASE | re.MULTILINE,
)
_HTML_SCRIPT_STYLE = re.compile(
    r"<(script|style)\b[^>]*>[\s\S]*?</\1\s*>",
    re.IGNORECASE,
)
_HTML_TABLE = re.compile(
    r"<table\b[^>]*>([\s\S]*?)</table\s*>",
    re.IGNORECASE,
)
_HTML_ROW = re.compile(
    r"<tr\b[^>]*>([\s\S]*?)</tr\s*>",
    re.IGNORECASE,
)
_HTML_CELL = re.compile(
    r"<(th|td)\b[^>]*>([\s\S]*?)</\1\s*>",
    re.IGNORECASE,
)
_HTML_TAG = re.compile(
    r"</?(?:html|body|div|span|table|thead|tbody|tfoot|tr|th|td|"
    r"script|style|p|br|h[1-6]|strong|b|em|i|u|ul|ol|li|"
    r"blockquote|section|article|header|footer|pre|code)\b[^>]*>",
    re.IGNORECASE,
)
_PLACEHOLDER = re.compile(
    r"(?:\bTBD\b|\bTODO\b|\bFIXME\b|\[insert\b|\[placeholder\]|lorem ipsum|example\.com)",
    re.IGNORECASE,
)
_WORD = re.compile(r"[A-Za-z0-9][A-Za-z0-9_-]{2,}")
_PREVIEW_MARKER = re.compile(
    r"Large\s+source\s+preserved\s+for\s+document\s+generation\s*:.*?hidden\s+in\s+(?:the\s+)?chat\s+preview",
    re.IGNORECASE | re.DOTALL,
)
_PRODUCTION_DIRECTIVES = re.compile(
    r"(?:document|artifact|pdf|source|generation|internal|final)\s+(?:production|creation|recovery)?\s*(?:requirements?|instructions?|directives?)",
    re.IGNORECASE,
)
_COMMAND_TITLE = re.compile(
    r"\b(?:add|create|make|generate|rename|convert|revise|update)\b.*\b(?:table|pdf|document|version|file)\b",
    re.IGNORECASE,
)
_MATH_WORD_CONCATENATION = re.compile(
    r"(?:slopeorrateofchange|initialvalueorverticalintercept|Importantlogarithmrulesare|Antiderivativeof|iscorrect)",
    re.IGNORECASE,
)
_MATH_TEXT_COMMAND = re.compile(
    r"\\(?:text|textrm|textsf|texttt|mathrm|mathbf|mathit|operatorname)\s*\{([^{}]*)\}",
    re.IGNORECASE,
)
_MATH_COMMAND = re.compile(r"\\[A-Za-z]+")
_MATH_SEMANTIC_COMMAND = re.compile(
    r"\\(?:frac|sqrt|sum|prod|int|lim|log|sin|cos|tan|exp|det|"
    r"lambda|mu|rho|sigma|theta|alpha|beta|gamma|delta|Delta|"
    r"pi|infty|partial|nabla|times|cdot|pm|le|ge|ne|approx|"
    r"begin|end|left|right)\b"
)
_MATH_RELATION = re.compile(r"(?:=|≤|≥|≠|≈|<|>|\\(?:le|ge|ne|approx)\b)")
_MATH_OPERATOR = re.compile(r"(?:[+*/^]|(?<![A-Za-z])-|√|∫|∑|∏|±|×|÷|\\(?:times|cdot|pm)\b)")
_MATH_STRUCTURED_VARIABLE = re.compile(
    r"(?:\b[A-Za-z]\s*(?:_[{A-Za-z0-9]+|_[A-Za-z0-9]+|\^[{A-Za-z0-9]+|\^[A-Za-z0-9]+)"
    r"|\\(?:lambda|mu|rho|sigma|theta|alpha|beta|gamma|delta|Delta)\s*(?:_[{A-Za-z0-9]+|_[A-Za-z0-9]+)?)"
)
_MATH_PROSE_CUE = re.compile(
    r"(?:[.!?](?:\s|$)|\b(?:because|therefore|however|which|while|although|"
    r"whenever|correct|shows|means|indicates|demonstrates|explains)\b)",
    re.IGNORECASE,
)
_MATH_NON_PROSE_WORDS = {
    "avg",
    "cos",
    "det",
    "dx",
    "dy",
    "exp",
    "log",
    "max",
    "mean",
    "min",
    "mod",
    "probability",
    "rate",
    "sin",
    "tan",
    "time",
    "total",
    "variance",
}
_INLINE_SUBSCRIPT_DIGITS = str.maketrans(
    "0123456789",
    "₀₁₂₃₄₅₆₇₈₉",
)


@dataclass(frozen=True, slots=True)
class ArtifactQualityIssue:
    code: str
    message: str
    severity: str = "warning"

    def to_dict(self) -> dict[str, str]:
        return {
            "code": self.code,
            "message": self.message,
            "severity": self.severity,
        }


@dataclass(slots=True)
class ArtifactQualityReport:
    status: str = "passed"
    page_or_slide_count: int = 0
    issues: list[ArtifactQualityIssue] = field(
        default_factory=list
    )

    @property
    def error_count(self) -> int:
        return sum(
            1
            for issue in self.issues
            if issue.severity == "error"
        )

    @property
    def warning_count(self) -> int:
        return sum(
            1
            for issue in self.issues
            if issue.severity == "warning"
        )

    def add(
        self,
        code: str,
        message: str,
        *,
        severity: str = "warning",
    ) -> None:
        if any(
            issue.code == code
            and issue.message == message
            for issue in self.issues
        ):
            return

        self.issues.append(
            ArtifactQualityIssue(
                code=code,
                message=message,
                severity=severity,
            )
        )
        self._refresh_status()

    def _refresh_status(self) -> None:
        if self.error_count:
            self.status = "failed"
        elif self.warning_count:
            self.status = "passed_with_warnings"
        else:
            self.status = "passed"

    def merge(
        self,
        other: "ArtifactQualityReport",
    ) -> "ArtifactQualityReport":
        for issue in other.issues:
            self.add(
                issue.code,
                issue.message,
                severity=issue.severity,
            )
        self.page_or_slide_count = max(
            self.page_or_slide_count,
            other.page_or_slide_count,
        )
        self._refresh_status()
        return self

    @classmethod
    def from_dict(
        cls,
        payload: dict[str, object],
    ) -> "ArtifactQualityReport":
        report = cls(
            page_or_slide_count=int(
                payload.get(
                    "page_or_slide_count",
                    0,
                )
            ),
        )
        raw_issues = payload.get("issues", [])
        if isinstance(raw_issues, list):
            for raw_issue in raw_issues:
                if not isinstance(raw_issue, dict):
                    continue
                report.add(
                    str(raw_issue.get("code", "unknown")),
                    str(
                        raw_issue.get(
                            "message",
                            "Artifact quality issue.",
                        )
                    ),
                    severity=str(
                        raw_issue.get(
                            "severity",
                            "warning",
                        )
                    ),
                )
        report._refresh_status()
        return report

    def to_dict(self) -> dict[str, object]:
        return {
            "status": self.status,
            "page_or_slide_count": self.page_or_slide_count,
            "error_count": self.error_count,
            "warning_count": self.warning_count,
            "issues": [
                issue.to_dict()
                for issue in self.issues
            ],
        }


def equation_expression_is_structurally_valid(
    expression: str,
) -> bool:
    """Return whether an equation block contains a real mathematical core.

    Mathematical publishing frequently uses alphabetic labels in ``\\text{}``
    or ``\\mathrm{}``, Greek command names, units, and aligned derivations. A
    raw word-count therefore rejects valid engineering equations. This check
    instead looks for relations, operators, variables, numbers, and semantic
    math commands, while still rejecting prose sentences accidentally wrapped
    in display-math delimiters.
    """

    raw = expression.strip().strip("$").strip()
    if not raw:
        return False

    label_fragments = _MATH_TEXT_COMMAND.findall(raw)
    math_core = _MATH_TEXT_COMMAND.sub(" ", raw)

    relation_count = len(_MATH_RELATION.findall(math_core))
    operator_count = len(_MATH_OPERATOR.findall(math_core))
    semantic_command_count = len(
        _MATH_SEMANTIC_COMMAND.findall(math_core)
    )
    structured_variable_count = len(
        _MATH_STRUCTURED_VARIABLE.findall(math_core)
    )
    numeric_count = len(
        re.findall(r"(?<![A-Za-z])\d+(?:[.,]\d+)?", math_core)
    )
    simple_variable_count = len(
        re.findall(r"(?<![A-Za-z\\])(?:[A-Za-z])(?![A-Za-z])", math_core)
    )

    residual = _MATH_COMMAND.sub(" ", math_core)
    residual = re.sub(r"[_^{}\[\](),;&|]", " ", residual)
    residual = re.sub(r"[=+*/<>≤≥≠≈±×÷√∫∑∏0-9.:%-]", " ", residual)
    prose_words = [
        word
        for word in re.findall(r"[A-Za-z]{2,}", residual)
        if word.casefold() not in _MATH_NON_PROSE_WORDS
    ]

    label_word_count = sum(
        len(re.findall(r"[A-Za-z]{2,}", fragment))
        for fragment in label_fragments
    )
    has_prose_cue = _MATH_PROSE_CUE.search(raw) is not None

    # A prose sentence with only a small inline calculation is not a display
    # equation. Keep it as prose so words never collapse in math rendering.
    if (
        len(prose_words) >= 6
        and has_prose_cue
        and relation_count <= 1
        and semantic_command_count == 0
        and structured_variable_count <= 1
    ):
        return False

    math_score = 0
    math_score += min(relation_count, 2) * 2
    math_score += min(operator_count, 3)
    math_score += min(semantic_command_count, 3)
    math_score += min(structured_variable_count, 2)
    math_score += 1 if numeric_count else 0
    math_score += 1 if simple_variable_count >= 2 else 0

    if math_score >= 2:
        return True

    # Text-only math commands are labels, not equations. They are valid only
    # when a mathematical core remains alongside them.
    if label_word_count >= 4 and math_score == 0:
        return False

    return len(prose_words) <= 2 and (
        relation_count > 0
        or operator_count > 0
        or semantic_command_count > 0
        or structured_variable_count > 0
    )


def clean_inline_markdown(value: str) -> str:
    """Return readable plain text without leaking Markdown tokens."""

    normalized = _HTML_SCRIPT_STYLE.sub("", value)
    normalized = _HTML_TAG.sub(" ", normalized)
    normalized = html_lib.unescape(normalized)
    normalized = _MARKDOWN_LINK.sub(
        lambda match: match.group(1),
        normalized,
    )
    normalized = _INLINE_MATH.sub(
        lambda match: match.group(1),
        normalized,
    )

    previous = None
    while previous != normalized:
        previous = normalized
        normalized = _MARKDOWN_TOKEN.sub(
            lambda match: match.group(1),
            normalized,
        )

    normalized = normalized.replace(r"\*", "*")
    normalized = normalized.replace(r"\_", "_")
    normalized = normalized.replace("**", "")
    normalized = normalized.replace("~~", "")
    replacements = {
        r"\alpha": "α",
        r"\beta": "β",
        r"\gamma": "γ",
        r"\delta": "δ",
        r"\Delta": "Δ",
        r"\lambda": "λ",
        r"\mu": "μ",
        r"\rho": "ρ",
        r"\sigma": "σ",
        r"\nabla": "∇",
        r"\pi": "π",
        r"\theta": "θ",
        r"\le": "≤",
        r"\ge": "≥",
        r"\ne": "≠",
        r"\pm": "±",
        r"\times": "×",
        r"\cdot": "·",
        r"\rightarrow": "→",
        r"\implies": "⇒",
    }
    for source, target in replacements.items():
        normalized = normalized.replace(source, target)
    normalized = re.sub(
        r"_(\d+)",
        lambda match: match.group(1).translate(
            _INLINE_SUBSCRIPT_DIGITS
        ),
        normalized,
    )
    return " ".join(normalized.split())


def _plain_html_text(value: str) -> str:
    without_blocks = _HTML_SCRIPT_STYLE.sub("", value)
    without_tags = re.sub(r"<[^>]+>", " ", without_blocks)
    return " ".join(
        html_lib.unescape(without_tags).split()
    ).strip()


def _html_table_to_markdown(match: re.Match[str]) -> str:
    rows: list[list[str]] = []
    for row_match in _HTML_ROW.finditer(match.group(1)):
        cells = [
            _plain_html_text(cell.group(2)).replace("|", "\\|")
            for cell in _HTML_CELL.finditer(row_match.group(1))
        ]
        if cells:
            rows.append(cells)
    if not rows:
        return _plain_html_text(match.group(0))

    width = max(len(row) for row in rows)
    normalized_rows = [
        row + [""] * (width - len(row))
        for row in rows
    ]
    header = normalized_rows[0]
    markdown = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join("---" for _ in header) + " |",
    ]
    markdown.extend(
        "| " + " | ".join(row) + " |"
        for row in normalized_rows[1:]
    )
    return "\n".join(markdown)


def _html_fragment_to_markdown(value: str) -> str:
    if _RAW_HTML.search(value) is None and not re.search(
        r"</?(?:p|br|h[1-6]|li|blockquote|section|article)\b",
        value,
        re.IGNORECASE,
    ):
        return value

    normalized = _HTML_SCRIPT_STYLE.sub("", value)
    normalized = _HTML_TABLE.sub(
        _html_table_to_markdown,
        normalized,
    )
    normalized = re.sub(
        r"<h([1-6])\b[^>]*>([\s\S]*?)</h\1\s*>",
        lambda match: (
            "\n\n"
            + "#" * int(match.group(1))
            + " "
            + _plain_html_text(match.group(2))
            + "\n\n"
        ),
        normalized,
        flags=re.IGNORECASE,
    )
    normalized = re.sub(
        r"<li\b[^>]*>([\s\S]*?)</li\s*>",
        lambda match: "\n- " + _plain_html_text(match.group(1)),
        normalized,
        flags=re.IGNORECASE,
    )
    normalized = re.sub(
        r"<br\s*/?>",
        "\n",
        normalized,
        flags=re.IGNORECASE,
    )
    normalized = re.sub(
        r"</?(?:p|div|section|article|header|footer|ul|ol|blockquote)\b[^>]*>",
        "\n\n",
        normalized,
        flags=re.IGNORECASE,
    )
    normalized = _HTML_TAG.sub(" ", normalized)
    normalized = html_lib.unescape(normalized)
    normalized = re.sub(r"[ \t]+\n", "\n", normalized)
    normalized = re.sub(r"\n[ \t]+", "\n", normalized)
    normalized = re.sub(r"\n{4,}", "\n\n\n", normalized)
    return normalized


def _normalize_html_outside_code_fences(value: str) -> str:
    output: list[str] = []
    buffer: list[str] = []
    in_fence = False

    def flush_buffer() -> None:
        if not buffer:
            return
        output.extend(
            _html_fragment_to_markdown(
                "\n".join(buffer)
            ).splitlines()
        )
        buffer.clear()

    for line in value.splitlines():
        if line.strip().startswith("```"):
            if in_fence:
                output.append(line.rstrip())
                in_fence = False
            else:
                flush_buffer()
                in_fence = True
                output.append(line.rstrip())
            continue
        if in_fence:
            output.append(line.rstrip())
        else:
            buffer.append(line)
    flush_buffer()
    return "\n".join(output)


def normalize_markdown_source(content: str) -> str:
    """Normalize provider output before parsing it into the document IR."""

    normalized = _INTERNAL_CONTEXT.sub("", content).strip()
    # Providers sometimes wrap an entire response in a ``markdown`` fence
    # while also emitting legitimate nested code fences.  A non-greedy regex
    # mistakes the first nested closing fence for the wrapper boundary.  Strip
    # the explicit first/last wrapper lines before processing inner fences.
    wrapper_lines = normalized.splitlines()
    outer_wrapper_close: int | None = None
    fence_stack: list[int] = []
    for line_index, wrapper_line in enumerate(wrapper_lines):
        fence = re.fullmatch(
            r"```([^`]*)",
            wrapper_line.strip(),
        )
        if fence is None:
            continue
        info = fence.group(1).strip()
        if info:
            fence_stack.append(line_index)
        elif fence_stack:
            opener = fence_stack.pop()
            if opener == 0:
                outer_wrapper_close = line_index

    if (
        len(wrapper_lines) >= 2
        and re.fullmatch(
            r"```(?:md|markdown)[ \t]*",
            wrapper_lines[0].strip(),
            re.IGNORECASE,
        )
        and outer_wrapper_close == len(wrapper_lines) - 1
    ):
        normalized = "\n".join(wrapper_lines[1:-1]).strip()

    outer_match = _OUTER_FENCE.match(normalized)
    if (
        outer_match is not None
        and "```" not in outer_match.group(1)
    ):
        normalized = outer_match.group(1).strip()

    normalized = _MARKDOWN_FENCE.sub(
        lambda match: match.group(1).strip(),
        normalized,
    )
    normalized = _normalize_html_outside_code_fences(
        normalized
    )

    lines: list[str] = []
    for raw_line in normalized.splitlines():
        if _HORIZONTAL_RULE.match(raw_line):
            lines.append("")
            continue
        lines.append(raw_line.rstrip())

    normalized = "\n".join(lines)
    normalized = re.sub(r"\n{4,}", "\n\n\n", normalized)
    return normalized.strip()


def _iter_text(document: ArtifactDocument) -> Iterable[str]:
    yield document.title
    if document.subtitle:
        yield document.subtitle

    for section in document.sections:
        yield section.title
        for block in section.blocks:
            if isinstance(block, ParagraphBlock):
                yield block.text
            elif isinstance(block, QuoteBlock):
                yield block.text
                if block.attribution:
                    yield block.attribution
            elif isinstance(block, CalloutBlock):
                yield block.title
                yield block.text
            elif isinstance(block, BulletListBlock):
                yield from block.items
            elif isinstance(block, TableBlock):
                if block.caption:
                    yield block.caption
                yield from block.columns
                for row in block.rows:
                    yield from row
            elif isinstance(block, ChartBlock):
                yield block.title
                yield from block.labels
                if block.caption:
                    yield block.caption
                for series in block.series:
                    yield series.name
            elif isinstance(block, CodeBlock):
                continue
            elif isinstance(block, EquationBlock):
                if block.label:
                    yield block.label
            elif isinstance(block, DiagramBlock):
                yield block.title
                yield from block.steps
            elif isinstance(block, PageBreakBlock):
                continue


def normalize_document_structure(
    document: ArtifactDocument,
) -> ArtifactDocument:
    """Clean every user-visible IR field while preserving literal code/math.

    Parsing is the primary normalization boundary, but artifacts can also be
    reconstructed from stored versions and transformed by layout/revision
    code.  Applying this final typed pass immediately before validation keeps
    create, revise, export, durable-source recovery, and attached-document
    redesign behavior consistent without corrupting code or equations.
    """

    def optional(value: str | None) -> str | None:
        if value is None:
            return None
        cleaned = clean_inline_markdown(value)
        return cleaned or None

    normalized_sections: list[ArtifactSection] = []
    for section in document.sections:
        normalized_blocks = []
        for block in section.blocks:
            if isinstance(block, ParagraphBlock):
                normalized_blocks.append(
                    replace(block, text=clean_inline_markdown(block.text))
                )
            elif isinstance(block, QuoteBlock):
                normalized_blocks.append(
                    replace(
                        block,
                        text=clean_inline_markdown(block.text),
                        attribution=optional(block.attribution),
                    )
                )
            elif isinstance(block, CalloutBlock):
                normalized_blocks.append(
                    replace(
                        block,
                        title=clean_inline_markdown(block.title),
                        text=clean_inline_markdown(block.text),
                    )
                )
            elif isinstance(block, BulletListBlock):
                normalized_blocks.append(
                    replace(
                        block,
                        items=tuple(
                            clean_inline_markdown(item)
                            for item in block.items
                        ),
                    )
                )
            elif isinstance(block, TableBlock):
                normalized_blocks.append(
                    replace(
                        block,
                        columns=tuple(
                            clean_inline_markdown(column)
                            for column in block.columns
                        ),
                        rows=tuple(
                            tuple(
                                clean_inline_markdown(cell)
                                for cell in row
                            )
                            for row in block.rows
                        ),
                        caption=optional(block.caption),
                    )
                )
            elif isinstance(block, ChartBlock):
                normalized_blocks.append(
                    replace(
                        block,
                        title=clean_inline_markdown(block.title),
                        labels=tuple(
                            clean_inline_markdown(label)
                            for label in block.labels
                        ),
                        series=tuple(
                            ChartSeries(
                                name=clean_inline_markdown(series.name),
                                values=series.values,
                            )
                            for series in block.series
                        ),
                        caption=optional(block.caption),
                        x_label=optional(block.x_label),
                        y_label=optional(block.y_label),
                    )
                )
            elif isinstance(block, DiagramBlock):
                normalized_blocks.append(
                    replace(
                        block,
                        title=clean_inline_markdown(block.title),
                        steps=tuple(
                            clean_inline_markdown(step)
                            for step in block.steps
                        ),
                    )
                )
            elif isinstance(block, EquationBlock):
                normalized_blocks.append(
                    replace(block, label=optional(block.label))
                )
            else:
                # Literal code and page-break metadata are intentionally not
                # treated as prose. Their syntax must survive rendering.
                normalized_blocks.append(block)

        normalized_sections.append(
            replace(
                section,
                title=clean_inline_markdown(section.title),
                blocks=tuple(normalized_blocks),
            )
        )

    return replace(
        document,
        title=clean_inline_markdown(document.title),
        subtitle=optional(document.subtitle),
        author=optional(document.author),
        sections=tuple(normalized_sections),
    )


def _source_keywords(
    source_snapshot: dict[str, Any] | None,
) -> set[str]:
    if not source_snapshot:
        return set()

    source_text = " ".join(
        str(source_snapshot.get(key, ""))
        for key in ("summary", "content")
    )
    stop_words = {
        "about",
        "after",
        "again",
        "also",
        "and",
        "are",
        "create",
        "document",
        "from",
        "have",
        "into",
        "make",
        "professional",
        "report",
        "that",
        "the",
        "their",
        "this",
        "using",
        "with",
        "your",
    }
    return {
        word.casefold()
        for word in _WORD.findall(source_text)
        if word.casefold() not in stop_words
    }


def validate_document_quality(
    document: ArtifactDocument,
    *,
    source_snapshot: dict[str, Any] | None = None,
) -> ArtifactQualityReport:
    report = ArtifactQualityReport()
    normalized_titles: set[str] = set()

    if not document.title.strip():
        report.add(
            "missing_title",
            "The artifact has no usable title.",
            severity="error",
        )
    elif _COMMAND_TITLE.search(document.title):
        report.add(
            "command_used_as_title",
            "A user action instruction was used as the document title.",
            severity="error",
        )

    for section in document.sections:
        key = section.title.casefold().strip()
        if not key:
            report.add(
                "blank_section_title",
                "A section has a blank title.",
                severity="error",
            )
        elif section.level <= 2 and key in normalized_titles:
            report.add(
                "duplicate_section",
                f"Duplicate section title: {section.title}",
            )
        elif section.level <= 2:
            normalized_titles.add(key)

        if not section.blocks:
            report.add(
                "empty_section",
                f"Section has no content: {section.title}",
                severity="error",
            )

        for block in section.blocks:
            if isinstance(block, ParagraphBlock):
                if len(block.text) > 8_000:
                    report.add(
                        "oversized_paragraph",
                        (
                            f"Section {section.title!r} contains "
                            "an oversized paragraph."
                        ),
                    )
            elif isinstance(block, BulletListBlock):
                if len(block.items) > 30:
                    report.add(
                        "oversized_list",
                        (
                            f"Section {section.title!r} contains "
                            "more than thirty list items."
                        ),
                    )
            elif isinstance(block, TableBlock):
                if len(block.columns) > 8:
                    report.add(
                        "wide_table",
                        (
                            f"Table in {section.title!r} has more "
                            "than eight columns and may be difficult "
                            "to render cleanly."
                        ),
                    )
                if any(
                    len(row) != len(block.columns)
                    for row in block.rows
                ):
                    report.add(
                        "invalid_table_shape",
                        (
                            f"Table in {section.title!r} has rows "
                            "that do not match its column count."
                        ),
                        severity="error",
                    )
            elif isinstance(block, CodeBlock):
                if "-----+" in block.code or "+-----" in block.code:
                    report.add(
                        "ascii_diagram",
                        (
                            f"Section {section.title!r} contains an "
                            "ASCII box diagram that should be converted "
                            "to a structured description."
                        ),
                    )
            elif isinstance(block, DiagramBlock):
                if len(block.steps) < 2:
                    report.add(
                        "underspecified_diagram",
                        f"Diagram in {section.title!r} has fewer than two steps.",
                    )
                if len(block.steps) > 12:
                    report.add(
                        "oversized_diagram",
                        f"Diagram in {section.title!r} has more than twelve steps.",
                    )
            elif isinstance(block, EquationBlock):
                if len(block.expression) > 1_000:
                    report.add(
                        "oversized_equation",
                        f"Equation in {section.title!r} is too long to render reliably.",
                    )
                if not equation_expression_is_structurally_valid(
                    block.expression
                ):
                    report.add(
                        "prose_classified_as_equation",
                        f"A prose sentence in {section.title!r} was classified as an equation.",
                        severity="error",
                    )

    combined_text = "\n".join(_iter_text(document))

    if _PREVIEW_MARKER.search(combined_text):
        report.add(
            "compact_preview_leaked_into_document",
            "A compact chat-preview marker leaked into the generated document.",
            severity="error",
        )
    if _PRODUCTION_DIRECTIVES.search(combined_text):
        report.add(
            "production_directives_rendered",
            "Internal document-production instructions were printed as document body content.",
            severity="error",
        )
    if _RAW_MARKER.search(combined_text):
        report.add(
            "raw_markdown_marker",
            "Raw Markdown markers remain in the structured document.",
            severity="error",
        )
    if _RAW_HTML.search(combined_text):
        report.add(
            "raw_html_marker",
            "Raw HTML markup remains in the structured document.",
            severity="error",
        )
    if _PLACEHOLDER.search(combined_text):
        report.add(
            "placeholder_content",
            "The artifact contains placeholder or unfinished content.",
            severity="error",
        )
    if len(combined_text.strip()) < 120:
        report.add(
            "insufficient_content",
            "The generated artifact contains too little useful content.",
            severity="error",
        )

    source_visualization_count = count_authentic_chart_blocks(
        str(source_snapshot.get("content", ""))
        if source_snapshot
        else None
    )
    document_visualization_count = sum(
        1
        for section in document.sections
        for block in section.blocks
        if isinstance(block, ChartBlock)
    )
    if document_visualization_count < source_visualization_count:
        report.add(
            "missing_source_visualization",
            (
                "One or more source charts were not preserved in the "
                "generated artifact."
            ),
            severity="error",
        )

    keywords = _source_keywords(source_snapshot)
    if keywords:
        document_words = {
            word.casefold()
            for word in _WORD.findall(combined_text)
        }
        overlap = keywords.intersection(document_words)
        minimum_overlap = 1 if len(keywords) < 4 else 2
        if len(overlap) < minimum_overlap:
            report.add(
                "source_topic_mismatch",
                (
                    "The generated document does not appear to preserve "
                    "the resolved source topic."
                ),
                severity="error",
            )

    return report


def _font_has_embedded_program(font_object: object) -> bool:
    try:
        font = font_object.get_object()
    except Exception:
        return False

    descriptor = font.get("/FontDescriptor")
    if descriptor is not None:
        try:
            descriptor = descriptor.get_object()
        except Exception:
            descriptor = None
        if descriptor is not None and any(
            key in descriptor
            for key in ("/FontFile", "/FontFile2", "/FontFile3")
        ):
            return True

    descendants = font.get("/DescendantFonts")
    if descendants:
        for descendant in descendants:
            if _font_has_embedded_program(descendant):
                return True
    return False


def _inspect_pdf(
    path: Path,
    report: ArtifactQualityReport,
) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    report.page_or_slide_count = len(reader.pages)
    page_texts = [
        page.extract_text() or ""
        for page in reader.pages
    ]

    if (
        artifact_settings.enforce_single_pdf_page_limit
        and report.page_or_slide_count
        > artifact_settings.maximum_single_pdf_pages
    ):
        report.add(
            "single_pdf_page_limit_exceeded",
            (
                "The generated PDF contains "
                f"{report.page_or_slide_count} pages, exceeding the "
                f"configured single-PDF limit of "
                f"{artifact_settings.maximum_single_pdf_pages}."
            ),
            severity="error",
        )

    try:
        outline = reader.outline
    except Exception:
        outline = []
    if report.page_or_slide_count >= 12 and not outline:
        report.add(
            "missing_pdf_outline",
            "A long PDF should include navigable outline bookmarks.",
        )

    seen_fonts: set[str] = set()
    unembedded_fonts: set[str] = set()
    for page in reader.pages:
        try:
            resources = page.get("/Resources")
            if resources is None:
                continue
            resources = resources.get_object()
            fonts = resources.get("/Font")
            if fonts is None:
                continue
            fonts = fonts.get_object()
            for font_reference in fonts.values():
                font = font_reference.get_object()
                name = str(font.get("/BaseFont", "Unknown"))
                if name in seen_fonts:
                    continue
                seen_fonts.add(name)
                base14 = {
                    "/Helvetica", "/Helvetica-Bold", "/Helvetica-Oblique",
                    "/Helvetica-BoldOblique", "/Times-Roman", "/Times-Bold",
                    "/Times-Italic", "/Times-BoldItalic", "/Courier",
                    "/Courier-Bold", "/Courier-Oblique", "/Courier-BoldOblique",
                    "/Symbol", "/ZapfDingbats",
                }
                if (
                    name not in base14
                    and not _font_has_embedded_program(font_reference)
                ):
                    unembedded_fonts.add(name)
        except Exception:
            continue
    if unembedded_fonts:
        report.add(
            "unembedded_pdf_fonts",
            (
                "One or more PDF fonts are not embedded: "
                + ", ".join(sorted(unembedded_fonts))
            ),
        )

    for index, page_text in enumerate(page_texts, start=1):
        if not page_text.strip():
            report.add(
                "blank_pdf_page",
                f"PDF page {index} contains no extractable text.",
            )
        elif len(page_text.strip()) < 24:
            report.add(
                "sparse_pdf_page",
                f"PDF page {index} contains very little text.",
            )

    return "\n".join(page_texts)


def _inspect_docx(
    path: Path,
    report: ArtifactQualityReport,
) -> str:
    from docx import Document

    document = Document(str(path))
    report.page_or_slide_count = max(1, len(document.sections))
    parts = [
        paragraph.text
        for paragraph in document.paragraphs
        if paragraph.text
    ]
    for table in document.tables:
        for row in table.rows:
            parts.extend(
                cell.text
                for cell in row.cells
                if cell.text
            )
    return "\n".join(parts)


def _inspect_pptx(
    path: Path,
    report: ArtifactQualityReport,
) -> str:
    from pptx import Presentation

    presentation = Presentation(str(path))
    report.page_or_slide_count = len(presentation.slides)
    slide_text: list[str] = []

    for slide_index, slide in enumerate(
        presentation.slides,
        start=1,
    ):
        current_slide: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text:
                current_slide.append(text)
        joined = "\n".join(current_slide)
        slide_text.append(joined)
        character_count = len(joined)
        if character_count > 1_400:
            report.add(
                "pptx_text_overflow_risk",
                (
                    f"Slide {slide_index} contains {character_count} "
                    "characters and may be too dense."
                ),
            )
        if not joined.strip():
            report.add(
                "blank_pptx_slide",
                f"Slide {slide_index} contains no readable text.",
            )

    return "\n".join(slide_text)



def _inspect_zip(
    path: Path,
    report: ArtifactQualityReport,
) -> str:
    from pypdf import PdfReader

    extracted_pages: list[str] = []

    try:
        archive = ZipFile(path, mode="r")
    except BadZipFile as error:
        raise ValueError("PDF bundle is not a readable ZIP archive.") from error

    with archive:
        names = archive.namelist()
        if "manifest.json" not in names:
            report.add(
                "missing_bundle_manifest",
                "The PDF bundle does not contain manifest.json.",
                severity="error",
            )

        unsafe_names = [
            name
            for name in names
            if name.startswith(("/", "\\"))
            or ".." in Path(name).parts
        ]
        if unsafe_names:
            report.add(
                "unsafe_bundle_path",
                "The PDF bundle contains an unsafe archive path.",
                severity="error",
            )

        pdf_names = sorted(
            name
            for name in names
            if name.lower().endswith(".pdf")
        )
        if len(pdf_names) < 2:
            report.add(
                "insufficient_bundle_volumes",
                "A multi-volume PDF bundle must contain at least two PDFs.",
                severity="error",
            )

        if "manifest.json" in names:
            try:
                manifest = json.loads(
                    archive.read("manifest.json").decode("utf-8")
                )
                declared = int(manifest.get("volume_count", 0))
                if declared != len(pdf_names):
                    report.add(
                        "bundle_manifest_mismatch",
                        "The PDF bundle manifest volume count does not match the archive.",
                        severity="error",
                    )
            except (UnicodeDecodeError, json.JSONDecodeError, TypeError, ValueError):
                report.add(
                    "invalid_bundle_manifest",
                    "The PDF bundle manifest is invalid.",
                    severity="error",
                )

        for pdf_name in pdf_names:
            reader = PdfReader(io.BytesIO(archive.read(pdf_name)))
            report.page_or_slide_count += len(reader.pages)
            for page in reader.pages:
                extracted_pages.append(page.extract_text() or "")

    return "\n".join(extracted_pages)

def inspect_rendered_file(
    path: Path,
    *,
    format: str,
) -> ArtifactQualityReport:
    report = ArtifactQualityReport()

    if not path.is_file() or path.stat().st_size <= 0:
        report.add(
            "missing_output",
            "The rendered artifact file is missing or empty.",
            severity="error",
        )
        return report

    extracted_text = ""
    try:
        if format == "pdf":
            extracted_text = _inspect_pdf(path, report)
        elif format == "docx":
            extracted_text = _inspect_docx(path, report)
        elif format == "pptx":
            extracted_text = _inspect_pptx(path, report)
        elif format == "zip":
            extracted_text = _inspect_zip(path, report)
        else:
            report.add(
                "unsupported_validation_format",
                f"No output validator is configured for {format!r}.",
                severity="error",
            )
    except Exception as error:
        report.add(
            "unreadable_output",
            f"The rendered file could not be reopened: {error}",
            severity="error",
        )
        return report

    if report.page_or_slide_count < 1:
        report.add(
            "empty_output_structure",
            "The rendered artifact contains no pages or slides.",
            severity="error",
        )
    if _PREVIEW_MARKER.search(extracted_text):
        report.add(
            "rendered_compact_preview_marker",
            "The rendered output contains a compact source-preview marker.",
            severity="error",
        )
    if _MATH_WORD_CONCATENATION.search(extracted_text):
        report.add(
            "rendered_math_word_concatenation",
            "Rendered equations contain concatenated prose words.",
            severity="error",
        )
    if _RAW_MARKER.search(extracted_text):
        report.add(
            "rendered_markdown_marker",
            (
                "The rendered output contains Markdown-like tokens; these "
                "may be intentional code or mathematical notation."
            ),
            severity="warning",
        )
    if _RAW_HTML.search(extracted_text):
        report.add(
            "rendered_html_marker",
            (
                "The rendered output contains HTML-like tokens; these may "
                "be intentional source-code examples."
            ),
            severity="warning",
        )
    if _PLACEHOLDER.search(extracted_text):
        report.add(
            "rendered_placeholder_content",
            "The rendered output contains placeholder content.",
            severity="error",
        )
    if not extracted_text.strip():
        report.add(
            "no_extractable_text",
            "The rendered output contains no extractable text.",
        )

    return report
